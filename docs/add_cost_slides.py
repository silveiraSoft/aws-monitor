#!/usr/bin/env python3
"""Agrega slides de costos y scripts de limpieza al PPTX de documentación."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

PPTX_PATH = os.path.join(os.path.dirname(__file__), "aws-monitor-dev-guide.pptx")

DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT  = RGBColor(0x58, 0x9B, 0xD6)
GREEN   = RGBColor(0x4E, 0xC9, 0x4E)
YELLOW  = RGBColor(0xFF, 0xD7, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xA0, 0xA0, 0xA0)
RED_L   = RGBColor(0xFF, 0x79, 0x79)


def get_layout(prs):
    return prs.slides[-1].slide_layout


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def tb(slide, text, left, top, width, height,
       size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def slide_costs(prs):
    slide = prs.slides.add_slide(get_layout(prs))
    add_bg(slide)
    # Limpiar placeholders heredados
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)

    tb(slide, "Estimacion de Costos AWS  -  us-east-1",
       0.3, 0.15, 9.2, 0.55, size=22, bold=True, color=ACCENT)
    tb(slide, "El riesgo de gasto es por USO, no por existencia de recursos.",
       0.3, 0.72, 9.2, 0.35, size=13, color=YELLOW, italic=True)

    rows = [
        ("Recurso",                 "En reposo",    "Por uso",                       "Nota"),
        ("CloudFront",              "~$0/mes",       "$0.0085 / 10K req",             ""),
        ("API Gateway",             "~$0/mes",       "$3.50 / millon llamadas",       ""),
        ("Lambda (2 funciones)",    "~$0/mes",       "$0.20 / millon invocaciones",   ""),
        ("S3 (2 buckets)",          "< $0.01/mes",   "$0.023/GB almacenado",          ""),
        ("Bedrock Agent (Haiku)",   "~$0/mes",       "~$0.0008 / 1K tokens",          "Principal riesgo"),
        ("CDKToolkit (bootstrap)",  "$0/mes",        "-",                             "IAM+SSM gratis"),
        ("TOTAL EN REPOSO",         "< $0.02/mes",   "",                              ""),
    ]

    col_x = [0.3, 3.2, 5.1, 7.3]
    col_w = [2.8, 1.8, 2.1, 2.3]
    row_h = 0.42

    for ri, row in enumerate(rows):
        y = 1.12 + ri * row_h
        for ci, val in enumerate(row):
            if ri == 0:
                c = ACCENT
            elif ri == len(rows) - 1:
                c = YELLOW
            elif ci == 1:
                c = GREEN
            elif ci == 3 and val:
                c = RED_L
            elif ci == 2:
                c = GRAY
            else:
                c = WHITE
            tb(slide, val, col_x[ci], y, col_w[ci], row_h - 0.05,
               size=12, bold=(ri == 0 or (ri == len(rows)-1 and ci < 2)), color=c)

    tb(slide,
       "Haiku: input $0.80/M tokens  |  output $4.00/M tokens  "
       " |  Produccion -> Claude 3.5 Sonnet (~$3/$15 por M tokens)",
       0.3, 5.25, 9.2, 0.4, size=10, color=GRAY, italic=True)


def slide_cleanup(prs):
    slide = prs.slides.add_slide(get_layout(prs))
    add_bg(slide)
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)

    tb(slide, "Scripts de Limpieza  -  Gestion de Costos",
       0.3, 0.15, 9.2, 0.55, size=22, bold=True, color=ACCENT)

    # --- cleanup_deploy.py ---
    tb(slide, "cleanup_deploy.py  <-- EJECUTAR AL FINAL DEL DIA",
       0.3, 0.82, 5.6, 0.38, size=13, bold=True, color=GREEN)
    tb(slide,
       "Elimina todos los recursos con potencial de costo:\n"
       "  Lambda (2 funciones)\n"
       "  API Gateway\n"
       "  CloudFront distribution\n"
       "  S3 buckets (vaciados + eliminados)\n"
       "  Bedrock Agent + alias live",
       0.3, 1.22, 5.5, 1.9, size=12, color=WHITE)
    tb(slide, "Uso:", 0.3, 3.18, 5.5, 0.28, size=12, bold=True, color=ACCENT)
    tb(slide,
       "python cleanup_deploy.py\n"
       "python cleanup_deploy.py --dry-run",
       0.3, 3.48, 5.5, 0.6, size=12, color=YELLOW, italic=True)

    # --- cleanup_bootstrap.py ---
    tb(slide, "cleanup_bootstrap.py  <-- SOLO LIMPIEZA TOTAL",
       6.0, 0.82, 3.7, 0.38, size=13, bold=True, color=GRAY)
    tb(slide,
       "Elimina el stack CDKToolkit:\n"
       "  Bucket S3 de assets CDK\n"
       "  4 IAM Roles bootstrap\n"
       "  SSM Parameter version\n\n"
       "Costo: $0/mes\n"
       "No es necesario ejecutarlo\n"
       "diariamente.",
       6.0, 1.22, 3.7, 2.5, size=12, color=GRAY)

    # --- Flujo diario ---
    tb(slide, "Flujo recomendado de desarrollo:",
       0.3, 4.22, 9.2, 0.32, size=13, bold=True, color=ACCENT)
    tb(slide,
       "Manana -> python validate_aws_access.py -> npm run deploy -> [usar el agente]",
       0.3, 4.57, 9.2, 0.32, size=12, color=WHITE)
    tb(slide,
       "Tarde  -> python cleanup_deploy.py -> [ambiente limpio, sin gasto acumulado]",
       0.3, 4.92, 9.2, 0.32, size=12, color=GREEN)
    tb(slide, "Para redesplegar: python validate_aws_access.py && npm run deploy",
       0.3, 5.27, 9.2, 0.32, size=11, color=GRAY, italic=True)


def slide_bootstrap(prs):
    slide = prs.slides.add_slide(get_layout(prs))
    add_bg(slide)
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)

    tb(slide, "CDK Bootstrap  -  Comportamiento y Costo",
       0.3, 0.15, 9.2, 0.55, size=22, bold=True, color=ACCENT)
    tb(slide, "npx cdk bootstrap aws://369595298303/us-east-1",
       0.3, 0.75, 9.2, 0.38, size=14, color=YELLOW, italic=True)

    items = [
        (True,  "Idempotente",
                "Ejecutarlo varias veces hace UPDATE, no crea recursos duplicados."),
        (True,  "Rollback automatico",
                "Si falla, CloudFormation revierte al estado anterior. Sin recursos huerfanos."),
        (True,  "Costo $0/mes",
                "IAM Roles y SSM Parameter son gratuitos. S3 solo cobra por contenido (~KB)."),
        (False, "Requiere iam:GetRole",
                "Sin este permiso el bootstrap falla. Agregar en IAM Console -> usuario -> Add permissions."),
        (False, "Requiere modelo habilitado",
                "Habilitar Claude 3.5 Haiku en Bedrock -> Model access (us-east-1) antes del deploy."),
        (False, "Inference Profile obligatorio en us-east-1",
                "Usar prefijo 'us.' -> us.anthropic.claude-3-5-haiku-20241022-v1:0"),
    ]

    for i, (ok, title, body) in enumerate(items):
        y = 1.22 + i * 0.67
        icon = "OK  " if ok else "!   "
        color = GREEN if ok else YELLOW
        tb(slide, icon + title, 0.3, y, 2.9, 0.38, size=13, bold=True, color=color)
        tb(slide, body,         3.3, y, 6.3, 0.55, size=12, color=WHITE)


def main():
    prs = Presentation(PPTX_PATH)
    n_before = len(prs.slides)
    print(f"PPTX abierto: {n_before} slides existentes")

    slide_costs(prs)
    slide_cleanup(prs)
    slide_bootstrap(prs)

    prs.save(PPTX_PATH)
    print(f"PPTX guardado: {len(prs.slides)} slides totales (+{len(prs.slides)-n_before} nuevas)")
    print(f"Archivo: {PPTX_PATH}")


if __name__ == "__main__":
    main()
